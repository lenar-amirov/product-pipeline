#!/usr/bin/env python3
"""
generate-pptx.py — конвертирует output/presentation.md в output/presentation.pptx

Usage:
    python3 generate-pptx.py <initiative-folder>
    python3 generate-pptx.py vk-tickets-vk-video
"""

import sys
import re
import os
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("❌ python-pptx не установлен. Запусти: pip install python-pptx")
    sys.exit(1)

# --- Тема оформления ---
BG_COLOR     = RGBColor(0x18, 0x18, 0x1A)  # тёмный фон
ACCENT_COLOR = RGBColor(0x00, 0x7A, 0xFF)  # синяя полоса
TITLE_COLOR  = RGBColor(0xFF, 0xFF, 0xFF)  # белый заголовок
BODY_COLOR   = RGBColor(0xCC, 0xCC, 0xCC)  # светло-серый текст
META_COLOR   = RGBColor(0x66, 0x66, 0x66)  # серый для вспомогательного

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def parse_slides(md_content: str) -> list[dict]:
    """Парсит markdown в список слайдов {title, body, notes}."""
    # Убираем шапку файла (# заголовок и > комментарии)
    content = re.sub(r'^# .+\n', '', md_content, count=1)
    content = re.sub(r'^>.*\n?', '', content, flags=re.MULTILINE)

    raw_slides = re.split(r'\n---\n', content)
    slides = []

    for raw in raw_slides:
        raw = raw.strip()
        if not raw or raw == '[Claude заполнит после запуска команды]':
            continue

        slide = {}

        # Заголовок: ## Слайд N: Название
        m = re.match(r'^##\s+(?:Слайд\s+\d+[:.]\s*)?(.+)$', raw, re.MULTILINE)
        if m:
            slide['title'] = m.group(1).strip()
        else:
            m = re.match(r'^#\s+(.+)$', raw, re.MULTILINE)
            slide['title'] = m.group(1).strip() if m else ''

        # Убираем строку заголовка из тела
        body = re.sub(r'^#{1,3} .+$', '', raw, count=1, flags=re.MULTILINE).strip()

        # Заметки спикера — строки начинающиеся с >
        notes_lines, body_lines = [], []
        for line in body.split('\n'):
            if line.startswith('> '):
                notes_lines.append(line[2:])
            elif line == '>':
                notes_lines.append('')
            else:
                body_lines.append(line)

        slide['body']  = '\n'.join(body_lines).strip()
        slide['notes'] = '\n'.join(notes_lines).strip()
        slides.append(slide)

    return slides


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR


def add_accent_line(slide):
    shape = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.fill.background()


def add_slide_number(slide, num, total):
    tb = slide.shapes.add_textbox(Inches(12.5), Inches(7.05), Inches(0.7), Pt(20))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{num}/{total}"
    run.font.size = Pt(10)
    run.font.color.rgb = META_COLOR


def add_title(slide, title: str):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.0), Inches(1.3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR


def add_body(slide, body: str):
    if not body:
        return

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.85), Inches(12.0), Inches(5.2))
    tf = tb.text_frame
    tf.word_wrap = True

    first_para = True
    for line in body.split('\n'):
        stripped = line.strip()
        if not stripped:
            continue

        p = tf.paragraphs[0] if first_para else tf.add_paragraph()
        first_para = False

        # Определяем уровень отступа
        if re.match(r'^\s{2,}[-*]', line):
            p.level = 1
            p.space_before = Pt(2)
            stripped = re.sub(r'^\s*[-*]\s+', '', stripped)
        elif re.match(r'^[-*•]', stripped):
            p.level = 0
            p.space_before = Pt(5)
            stripped = re.sub(r'^[-*•]\s+', '', stripped)
        else:
            p.level = 0
            p.space_before = Pt(7)

        # Разбиваем по **жирному**
        parts = re.split(r'\*\*(.*?)\*\*', stripped)
        for i, part in enumerate(parts):
            if not part:
                continue
            part_clean = re.sub(r'[`_]', '', part)
            run = p.add_run()
            run.text = part_clean
            run.font.size = Pt(18)
            run.font.color.rgb = TITLE_COLOR if i % 2 == 1 else BODY_COLOR
            if i % 2 == 1:
                run.font.bold = True


def add_speaker_notes(slide, notes: str):
    if not notes:
        return
    slide.notes_slide.notes_text_frame.text = notes


def build_pptx(slides_data: list[dict], out_path: Path):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]
    total = len(slides_data)

    for i, data in enumerate(slides_data, 1):
        slide = prs.slides.add_slide(blank_layout)
        set_bg(slide)
        add_accent_line(slide)
        add_slide_number(slide, i, total)
        if data.get('title'):
            add_title(slide, data['title'])
        add_body(slide, data.get('body', ''))
        add_speaker_notes(slide, data.get('notes', ''))

    prs.save(str(out_path))


def main():
    if len(sys.argv) < 2:
        print("Usage: python3 generate-pptx.py <initiative-folder>")
        sys.exit(1)

    initiative = sys.argv[1]
    base = Path(__file__).parent

    md_path  = base / initiative / "output" / "presentation.md"
    out_path = base / initiative / "output" / "presentation.pptx"

    if not md_path.exists():
        print(f"❌ Файл не найден: {md_path}")
        sys.exit(1)

    slides_data = parse_slides(md_path.read_text(encoding='utf-8'))

    if not slides_data:
        print("❌ Слайды не найдены. Проверь формат presentation.md")
        sys.exit(1)

    build_pptx(slides_data, out_path)
    print(f"✅ Готово: {out_path}")
    print(f"   Слайдов: {len(slides_data)}")


if __name__ == '__main__':
    main()
